"""
Educational SaaS Platform - Comprehensive Presentation Generator

This script generates a professional PowerPoint presentation showcasing
the Educational SaaS Platform's features, architecture, and capabilities.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime


class PresentationGenerator:
    """Generate comprehensive Educational SaaS Platform presentation"""
    
    # Professional color scheme
    PRIMARY_COLOR = RGBColor(26, 35, 126)      # Deep Blue
    SECONDARY_COLOR = RGBColor(67, 160, 71)    # Green
    ACCENT_COLOR = RGBColor(255, 152, 0)       # Orange
    TEXT_COLOR = RGBColor(33, 33, 33)          # Dark Gray
    LIGHT_BG = RGBColor(245, 245, 245)         # Light Gray
    WHITE = RGBColor(255, 255, 255)
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
    def add_title_slide(self):
        """Slide 1: Title Slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.PRIMARY_COLOR
        
        # Platform Name
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Educational SaaS Platform"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = self.WHITE
        title_para.alignment = PP_ALIGN.CENTER
        
        # Tagline
        tagline_box = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(8), Inches(0.8))
        tagline_frame = tagline_box.text_frame
        tagline_frame.text = "Multi-Tenant B2B2C Educational Management Platform"
        tagline_para = tagline_frame.paragraphs[0]
        tagline_para.font.size = Pt(24)
        tagline_para.font.color.rgb = self.ACCENT_COLOR
        tagline_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "AI-Powered Learning • Blockchain Credentials • Mobile-First"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.color.rgb = self.WHITE
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Date
        date_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.4))
        date_frame = date_box.text_frame
        date_frame.text = f"Generated: {datetime.now().strftime('%B %Y')}"
        date_para = date_frame.paragraphs[0]
        date_para.font.size = Pt(14)
        date_para.font.color.rgb = self.WHITE
        date_para.alignment = PP_ALIGN.CENTER
        
    def add_executive_summary(self):
        """Slide 2: Executive Summary"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        self._add_slide_title(slide, "Executive Summary")
        
        # Content box
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        # Key value propositions
        value_props = [
            ("🎯 Complete Educational Ecosystem", 
             "End-to-end platform covering student management, learning, assessment, and analytics"),
            ("🏢 Multi-Tenant B2B2C Architecture", 
             "Scalable SaaS solution supporting multiple institutions with isolated data"),
            ("🤖 AI-Powered Intelligence", 
             "Smart recommendations, weakness detection, automated grading, and study assistance"),
            ("📱 Mobile-First Design", 
             "Native mobile apps built with Expo Router and React Native"),
            ("🔐 Enterprise Security", 
             "JWT authentication, RBAC, audit logging, and blockchain credentials"),
            ("☁️ Cloud-Native Infrastructure", 
             "Docker/AWS deployment with auto-scaling and high availability"),
        ]
        
        for i, (title, desc) in enumerate(value_props):
            p = tf.add_paragraph()
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.PRIMARY_COLOR
            p.space_after = Pt(4)
            
            p = tf.add_paragraph()
            p.text = desc
            p.font.size = Pt(12)
            p.font.color.rgb = self.TEXT_COLOR
            p.space_after = Pt(12)
            p.level = 1
            
    def add_platform_overview(self):
        """Slide 3: Platform Overview"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_slide_title(slide, "Platform Overview: Multi-Tenant B2B2C Architecture")
        
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        sections = [
            ("Architecture Model", [
                "B2B2C SaaS platform serving multiple educational institutions",
                "Complete data isolation with institution-based tenancy",
                "Row-level security (RLS) for data protection",
                "Shared infrastructure with isolated data stores"
            ]),
            ("Core Capabilities", [
                "Institution management with subscription billing",
                "User management with role-based access control (RBAC)",
                "Comprehensive audit logging for compliance",
                "Multi-role support: Students, Teachers, Parents, Admins, Super Admins"
            ]),
            ("Scalability Features", [
                "Horizontal scaling with load balancing",
                "Redis caching for performance optimization",
                "Celery distributed task queuing",
                "WebSocket support for real-time features"
            ])
        ]
        
        for section_title, items in sections:
            p = tf.add_paragraph()
            p.text = section_title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.PRIMARY_COLOR
            p.space_after = Pt(6)
            
            for item in items:
                p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(12)
                p.font.color.rgb = self.TEXT_COLOR
                p.level = 1
                p.space_after = Pt(4)
            
            # Add spacing between sections
            p = tf.add_paragraph()
            p.space_after = Pt(8)
            
    def add_technology_stack(self):
        """Slide 4: Technology Stack"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_slide_title(slide, "Technology Stack")
        
        # Create two columns
        left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(4.2), Inches(5.5))
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.2), Inches(5.5))
        
        # Backend Stack
        left_tf = left_box.text_frame
        left_tf.word_wrap = True
        
        p = left_tf.add_paragraph()
        p.text = "Backend Stack"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.PRIMARY_COLOR
        p.space_after = Pt(10)
        
        backend_items = [
            ("FastAPI", "High-performance async Python framework"),
            ("MySQL", "Relational database with SQLAlchemy ORM"),
            ("Redis", "In-memory cache and session store"),
            ("Celery", "Distributed task queue"),
            ("Alembic", "Database migration management"),
            ("Pydantic", "Data validation and settings"),
            ("PyTorch", "Deep learning framework"),
            ("Scikit-learn", "Machine learning library"),
            ("Sentence Transformers", "NLP embeddings"),
        ]
        
        for tech, desc in backend_items:
            p = left_tf.add_paragraph()
            p.text = f"{tech}: "
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.SECONDARY_COLOR
            run = p.runs[0]
            
            run2 = p.add_run()
            run2.text = desc
            run2.font.bold = False
            run2.font.color.rgb = self.TEXT_COLOR
            p.space_after = Pt(6)
        
        # Frontend & Mobile Stack
        right_tf = right_box.text_frame
        right_tf.word_wrap = True
        
        p = right_tf.add_paragraph()
        p.text = "Frontend & Mobile"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.PRIMARY_COLOR
        p.space_after = Pt(10)
        
        frontend_items = [
            ("React Native", "Cross-platform mobile framework"),
            ("Expo Router", "File-based navigation"),
            ("TypeScript", "Type-safe development"),
            ("Redux Toolkit", "State management"),
            ("React", "Web frontend framework"),
            ("Axios", "HTTP client"),
            ("WebSocket", "Real-time communication"),
        ]
        
        for tech, desc in frontend_items:
            p = right_tf.add_paragraph()
            p.text = f"{tech}: "
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.SECONDARY_COLOR
            run = p.runs[0]
            
            run2 = p.add_run()
            run2.text = desc
            run2.font.bold = False
            run2.font.color.rgb = self.TEXT_COLOR
            p.space_after = Pt(6)
            
    def add_core_features_by_role(self):
        """Slide 5: Core Features by User Role"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_slide_title(slide, "Core Features by User Role")
        
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        roles = [
            ("👨‍🎓 Students", [
                "Course enrollment and content access",
                "Assignment submission and tracking",
                "Exam participation with auto-grading",
                "AI Study Buddy and homework scanner",
                "Performance analytics and weakness detection",
                "Peer tutoring and collaboration tools"
            ]),
            ("👨‍🏫 Teachers", [
                "Course creation and content management",
                "Assignment and exam creation with rubrics",
                "Automated and manual grading tools",
                "Attendance tracking and reporting",
                "Student progress monitoring",
                "Resource library and content marketplace"
            ]),
            ("👨‍👩‍👧‍👦 Parents", [
                "Children's academic progress tracking",
                "Attendance and behavior monitoring",
                "Communication with teachers",
                "Payment and fee management",
                "Event and calendar access"
            ]),
            ("⚙️ Admins", [
                "Institution configuration and settings",
                "User management and role assignment",
                "Subscription and billing management",
                "Reports and analytics dashboards",
                "System health monitoring"
            ]),
            ("🔧 Super Admins", [
                "Multi-institution management",
                "Platform-wide configuration",
                "System monitoring and maintenance",
                "Cross-tenant reporting and analytics"
            ])
        ]
        
        for role_name, features in roles:
            p = tf.add_paragraph()
            p.text = role_name
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.PRIMARY_COLOR
            p.space_after = Pt(4)
            
            p = tf.add_paragraph()
            p.text = " • ".join(features)
            p.font.size = Pt(10)
            p.font.color.rgb = self.TEXT_COLOR
            p.level = 1
            p.space_after = Pt(10)
            
    def add_ai_ml_capabilities(self):
        """Slide 6: AI/ML Capabilities"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_slide_title(slide, "AI/ML Capabilities & Intelligent Features")
        
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        ai_features = [
            ("🤖 AI Study Buddy", [
                "GPT-4 powered conversational tutoring assistant",
                "Personalized study plan generation",
                "Context-aware learning recommendations",
                "Motivational insights and progress tracking"
            ]),
            ("📸 Smart Homework Scanner", [
                "OCR-based text extraction from images (Tesseract)",
                "Automatic math problem solving (SymPy)",
                "Step-by-step solution explanations",
                "AI-generated feedback and hints"
            ]),
            ("📊 Weakness Detection Engine", [
                "Chapter-wise performance analysis",
                "Spaced repetition system (SM-2 algorithm)",
                "Focus area prioritization with urgency scoring",
                "Personalized insight generation"
            ]),
            ("🎯 Smart Recommendation Engine", [
                "Collaborative filtering using cosine similarity",
                "Content effectiveness scoring and tracking",
                "Difficulty level auto-detection",
                "VARK learning style adaptation (Visual, Auditory, Reading, Kinesthetic)",
                "Personalized study path sequencing"
            ]),
            ("🔮 Predictive Analytics", [
                "Student performance prediction models",
                "At-risk student identification",
                "Exam topic prediction using ML",
                "Dropout risk analysis"
            ]),
            ("📚 Content Intelligence", [
                "Automatic content categorization",
                "Semantic search with embeddings",
                "Plagiarism detection",
                "Question bank generation"
            ])
        ]
        
        for feature_name, capabilities in ai_features:
            p = tf.add_paragraph()
            p.text = feature_name
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.ACCENT_COLOR
            p.space_after = Pt(4)
            
            for capability in capabilities:
                p = tf.add_paragraph()
                p.text = capability
                p.font.size = Pt(10)
                p.font.color.rgb = self.TEXT_COLOR
                p.level = 1
                p.space_after = Pt(2)
            
            p = tf.add_paragraph()
            p.space_after = Pt(6)
            
    def add_mobile_app_highlights(self):
        """Slide 7: Mobile App Highlights"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_slide_title(slide, "Mobile App: Expo Router Architecture")
        
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        sections = [
            ("📱 Cross-Platform Excellence", [
                "React Native with Expo for iOS, Android, and Web",
                "File-based routing with Expo Router",
                "TypeScript for type safety",
                "Bundle size optimized < 2MB for web"
            ]),
            ("🔐 Security & Authentication", [
                "Biometric authentication (Face ID/Touch ID)",
                "Secure token storage with Expo Secure Store",
                "Auto-refresh token mechanism",
                "OTP-based login support"
            ]),
            ("⚡ Performance Optimization", [
                "Lazy loading for heavy screens (AI Predictions, Scanner)",
                "Aggressive tree-shaking and code splitting",
                "AsyncStorage for web, SecureStore for native",
                "Dynamic imports for bundle optimization"
            ]),
            ("✨ Key Features", [
                "Role-based navigation (Student/Parent/Teacher)",
                "Deep linking support",
                "Real-time notifications with WebSocket",
                "Offline functionality with local caching",
                "Homework scanner with camera integration",
                "AI chat interface for study assistance"
            ]),
            ("🎨 User Experience", [
                "Native UI components and animations",
                "Responsive design for tablets and phones",
                "Dark mode support",
                "Accessibility features (screen reader support)"
            ])
        ]
        
        for section_title, items in sections:
            p = tf.add_paragraph()
            p.text = section_title
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.PRIMARY_COLOR
            p.space_after = Pt(5)
            
            for item in items:
                p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(11)
                p.font.color.rgb = self.TEXT_COLOR
                p.level = 1
                p.space_after = Pt(3)
            
            p = tf.add_paragraph()
            p.space_after = Pt(8)
            
    def add_security_compliance(self):
        """Slide 8: Security & Compliance"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_slide_title(slide, "Security & Compliance")
        
        # Create two columns
        left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(4.2), Inches(5.5))
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.2), Inches(5.5))
        
        # Authentication & Authorization
        left_tf = left_box.text_frame
        left_tf.word_wrap = True
        
        p = left_tf.add_paragraph()
        p.text = "Authentication & Authorization"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.PRIMARY_COLOR
        p.space_after = Pt(8)
        
        auth_items = [
            "JWT token-based authentication",
            "Bcrypt password hashing",
            "Role-Based Access Control (RBAC)",
            "18+ granular permissions",
            "5 system roles with hierarchy",
            "Session management with Redis",
            "OTP and biometric login",
            "Password reset workflows"
        ]
        
        for item in auth_items:
            p = left_tf.add_paragraph()
            p.text = f"✓ {item}"
            p.font.size = Pt(11)
            p.font.color.rgb = self.TEXT_COLOR
            p.space_after = Pt(4)
        
        # Data Protection & Compliance
        right_tf = right_box.text_frame
        right_tf.word_wrap = True
        
        p = right_tf.add_paragraph()
        p.text = "Data Protection & Compliance"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.PRIMARY_COLOR
        p.space_after = Pt(8)
        
        data_items = [
            "Row-Level Security (RLS) policies",
            "Institution-based data isolation",
            "Comprehensive audit logging",
            "Encrypted data at rest and in transit",
            "HTTPS/TLS for all communications",
            "Regular security scanning (Bandit)",
            "SQL injection prevention (ORM)",
            "XSS and CSRF protection",
            "Rate limiting and DDoS protection",
            "Backup and disaster recovery",
            "GDPR compliance ready",
            "SOC 2 preparation support"
        ]
        
        for item in data_items:
            p = right_tf.add_paragraph()
            p.text = f"✓ {item}"
            p.font.size = Pt(11)
            p.font.color.rgb = self.TEXT_COLOR
            p.space_after = Pt(4)
            
    def add_deployment_infrastructure(self):
        """Slide 9: Deployment & Infrastructure"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_slide_title(slide, "Deployment & Infrastructure")
        
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        sections = [
            ("🐳 Containerization", [
                "Docker containers for all services",
                "Docker Compose for local development",
                "Multi-stage builds for optimization",
                "Production and development Dockerfiles"
            ]),
            ("☁️ AWS Cloud Infrastructure", [
                "ECS/EC2 for application hosting",
                "RDS MySQL for database (Multi-AZ)",
                "ElastiCache Redis for caching",
                "S3 for file storage and backups",
                "CloudFront CDN for content delivery",
                "Application Load Balancer (ALB)",
                "Route 53 for DNS management",
                "CloudWatch for monitoring and logs"
            ]),
            ("🔄 CI/CD Pipeline", [
                "GitLab CI/CD automation",
                "Automated testing on every commit",
                "Code quality checks (Black, Ruff, MyPy)",
                "Security scanning (Bandit)",
                "Docker image building and registry",
                "Blue-green deployment strategy",
                "Automated rollback capabilities"
            ]),
            ("📊 Monitoring & Observability", [
                "Sentry for error tracking",
                "CloudWatch metrics and alarms",
                "Application performance monitoring",
                "Database query optimization",
                "Real-time health checks",
                "Automated alerting system"
            ]),
            ("🚀 Scalability", [
                "Horizontal auto-scaling",
                "Load balancing across instances",
                "Database connection pooling",
                "Redis caching layer",
                "Celery worker scaling",
                "CDN for static assets"
            ])
        ]
        
        for section_title, items in sections:
            p = tf.add_paragraph()
            p.text = section_title
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.PRIMARY_COLOR
            p.space_after = Pt(4)
            
            for item in items:
                p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(10)
                p.font.color.rgb = self.TEXT_COLOR
                p.level = 1
                p.space_after = Pt(2)
            
            p = tf.add_paragraph()
            p.space_after = Pt(6)
            
    def add_unique_differentiators(self):
        """Slide 10: Unique Differentiators"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_slide_title(slide, "Unique Differentiators & Innovation")
        
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        differentiators = [
            ("🔗 Blockchain Digital Credentials", [
                "Hyperledger Fabric network integration",
                "Tamper-proof certificate issuance",
                "Verifiable academic credentials",
                "Credential revocation support",
                "Decentralized verification system",
                "Complete credential history tracking"
            ]),
            ("🎓 Advanced Learning Features", [
                "Reverse classroom implementation",
                "Peer tutoring marketplace",
                "Learning style detection (VARK model)",
                "Gamification with achievement system",
                "Community service hour tracking",
                "Volunteer hours management"
            ]),
            ("📚 Content Ecosystem", [
                "Content marketplace for teachers",
                "External library integration (Khan Academy, Coursera, edX)",
                "Document vault with version control",
                "Resource sharing across institutions",
                "Collaborative content creation"
            ]),
            ("👥 Community & Engagement", [
                "Live events and webinar hosting",
                "Student yearbook digital platform",
                "Peer recognition system",
                "Goal setting with gamification",
                "Carpool coordination",
                "Merchandise store integration"
            ]),
            ("📊 Advanced Analytics", [
                "Institution health monitoring dashboard",
                "Predictive analytics for student success",
                "Performance trend analysis",
                "Custom report builder",
                "Data export and integration APIs",
                "Real-time analytics dashboards"
            ]),
            ("🔄 Seamless Integrations", [
                "REST API with comprehensive documentation",
                "Webhook support for external systems",
                "LMS integration capabilities",
                "Payment gateway integration (Stripe/PayPal)",
                "Email service integration (SendGrid)",
                "SMS notifications support"
            ])
        ]
        
        for feature_name, items in differentiators:
            p = tf.add_paragraph()
            p.text = feature_name
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.ACCENT_COLOR
            p.space_after = Pt(4)
            
            for item in items:
                p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(10)
                p.font.color.rgb = self.TEXT_COLOR
                p.level = 1
                p.space_after = Pt(2)
            
            p = tf.add_paragraph()
            p.space_after = Pt(6)
            
    def add_roadmap(self):
        """Slide 11: Roadmap & Future Enhancements"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_slide_title(slide, "Roadmap & Future Enhancements")
        
        content_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(8.2), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        roadmap = [
            ("Q2 2026 - Enhanced AI Capabilities", [
                "Advanced NLP for essay grading",
                "Voice-to-text assignment submission",
                "AI-powered proctoring for online exams",
                "Emotion recognition for engagement tracking",
                "Personalized learning path optimization"
            ]),
            ("Q3 2026 - Extended Platform Features", [
                "Video conferencing integration (Zoom/Teams)",
                "Virtual classroom environments",
                "Augmented Reality (AR) learning modules",
                "Advanced plagiarism detection with AI",
                "Multi-language support and localization"
            ]),
            ("Q4 2026 - Analytics & Insights", [
                "Advanced predictive analytics dashboard",
                "Parent engagement scoring",
                "Teacher effectiveness metrics",
                "Custom report builder with drag-and-drop",
                "Data warehouse for historical analysis"
            ]),
            ("2027 - Ecosystem Expansion", [
                "Marketplace for third-party plugins",
                "White-label solution for institutions",
                "API marketplace for developers",
                "Mobile SDK for custom app development",
                "Integration with national education standards"
            ]),
            ("Continuous Improvements", [
                "Performance optimization and scaling",
                "Enhanced security features",
                "UI/UX improvements based on user feedback",
                "Accessibility enhancements (WCAG 2.1 AA)",
                "Regular feature updates and bug fixes"
            ])
        ]
        
        for quarter, items in roadmap:
            p = tf.add_paragraph()
            p.text = quarter
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.PRIMARY_COLOR
            p.space_after = Pt(5)
            
            for item in items:
                p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(11)
                p.font.color.rgb = self.TEXT_COLOR
                p.level = 1
                p.space_after = Pt(3)
            
            p = tf.add_paragraph()
            p.space_after = Pt(8)
            
    def add_section_divider(self, title):
        """Add a section divider slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.SECONDARY_COLOR
        
        # Section Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = self.WHITE
        title_para.alignment = PP_ALIGN.CENTER
        
    def _add_slide_title(self, slide, title_text):
        """Add a formatted title to a slide"""
        # White background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.WHITE
        
        # Title bar
        title_bar = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(10), Inches(1)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.PRIMARY_COLOR
        title_bar.line.fill.background()
        
        # Title text
        title_frame = title_bar.text_frame
        title_frame.text = title_text
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self.WHITE
        title_para.alignment = PP_ALIGN.CENTER
        
    def generate(self, filename="Educational_SaaS_Platform_Presentation.pptx"):
        """Generate the complete presentation"""
        print("Generating Educational SaaS Platform Presentation...")
        
        # Add all slides
        print("  [1/13] Adding title slide...")
        self.add_title_slide()
        
        print("  [2/13] Adding executive summary...")
        self.add_executive_summary()
        
        print("  [3/13] Adding platform overview...")
        self.add_section_divider("Platform Architecture")
        self.add_platform_overview()
        
        print("  [4/13] Adding technology stack...")
        self.add_technology_stack()
        
        print("  [5/13] Adding core features...")
        self.add_section_divider("Features & Capabilities")
        self.add_core_features_by_role()
        
        print("  [6/13] Adding AI/ML capabilities...")
        self.add_ai_ml_capabilities()
        
        print("  [7/13] Adding mobile app highlights...")
        self.add_section_divider("Mobile & User Experience")
        self.add_mobile_app_highlights()
        
        print("  [8/13] Adding security & compliance...")
        self.add_section_divider("Security & Infrastructure")
        self.add_security_compliance()
        
        print("  [9/13] Adding deployment & infrastructure...")
        self.add_deployment_infrastructure()
        
        print("  [10/13] Adding unique differentiators...")
        self.add_section_divider("Innovation & Differentiation")
        self.add_unique_differentiators()
        
        print("  [11/13] Adding roadmap...")
        self.add_section_divider("Future Vision")
        self.add_roadmap()
        
        # Save presentation
        print(f"\n  Saving presentation to: {filename}")
        self.prs.save(filename)
        print(f"\n✓ Presentation generated successfully!")
        print(f"  Total slides: {len(self.prs.slides)}")
        print(f"  File: {filename}")
        

def main():
    """Main entry point"""
    generator = PresentationGenerator()
    generator.generate()


if __name__ == "__main__":
    main()
